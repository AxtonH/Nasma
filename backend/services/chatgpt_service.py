from openai import OpenAI
from backend.config.settings import Config
from datetime import datetime

def debug_log(message: str, category: str = "general"):
    """Conditional debug logging based on configuration"""
    if category == "odoo_data" and Config.DEBUG_ODOO_DATA:
        print(f"DEBUG: {message}")
    elif category == "bot_logic" and Config.DEBUG_BOT_LOGIC:
        print(f"DEBUG: {message}")
    elif category == "knowledge_base" and Config.DEBUG_KNOWLEDGE_BASE:
        print(f"DEBUG: {message}")
    elif category == "general" and Config.VERBOSE_LOGS:
        print(f"DEBUG: {message}")
import time
import json
import os
import threading

class ChatGPTService:
    """Service class for handling ChatGPT interactions using the latest available GPT model"""
    
    def __init__(self):
        self.client = None  # Initialize lazily to avoid hanging
        self.model = Config.GPT_MODEL  # Using configured GPT model
        self.fallback_model = getattr(Config, 'GPT_FALLBACK_MODEL', None)
        self.conversation_history = {}  # Store conversation history by thread_id
        self.storage_dir = "conversation_storage"
        self.lock = threading.Lock()  # Thread safety for file operations
        
        # Service dependencies (set via set_services)
        self.timeoff_service = None
        self.session_manager = None
        self.halfday_service = None
        self.reimbursement_service = None
        
        # Disable conversation storage
        # if not os.path.exists(self.storage_dir):
        #     os.makedirs(self.storage_dir)
    
    def set_services(self, timeoff_service, session_manager, halfday_service=None, reimbursement_service=None):
        """Wire external services for advanced functionality"""
        self.timeoff_service = timeoff_service
        self.session_manager = session_manager
        self.halfday_service = halfday_service
        self.reimbursement_service = reimbursement_service
    
    def _get_client(self):
        """Get OpenAI client, initializing lazily if needed"""
        if self.client is None:
            try:
                self.client = OpenAI(api_key=Config.OPENAI_API_KEY)
            except Exception as e:
                debug_log(f"Failed to initialize OpenAI client: {e}", "general")
                return None
        return self.client
    
    def _get_storage_file(self, thread_id):
        """Get the storage file path for a thread"""
        return os.path.join(self.storage_dir, f"thread_{thread_id}.json")

    def _get_summary_file(self, thread_id):
        return os.path.join(self.storage_dir, f"thread_{thread_id}_summary.json")

    def _load_summary(self, thread_id):
        try:
            path = self._get_summary_file(thread_id)
            if os.path.exists(path):
                with open(path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    return data.get('summary', '')
        except Exception as e:
            debug_log(f"Error loading summary: {e}", "general")
        return ''

    def _save_summary(self, thread_id, summary_text):
        try:
            path = self._get_summary_file(thread_id)
            with open(path, 'w', encoding='utf-8') as f:
                json.dump({'summary': summary_text or ''}, f, ensure_ascii=False, indent=2)
        except Exception as e:
            debug_log(f"Error saving summary: {e}", "general")
    
    def _load_conversation_history(self, thread_id):
        """Load conversation history from persistent storage - DISABLED"""
        # Storage disabled - return empty history
        return []
    
    def _save_conversation_history(self, thread_id, history):
        """Save conversation history to persistent storage - DISABLED"""
        # Storage disabled - no longer saving conversation history
        pass
    
    def get_response(self, message, thread_id, employee_data=None):
        """Get response from configured GPT model using Chat Completions API"""
        try:
            debug_log(f"Getting {self.model} response for message: '{message[:100]}...'", "bot_logic")
            debug_log(f"Thread ID provided: {thread_id}", "bot_logic")
            debug_log(f"Employee data provided: {employee_data is not None}", "bot_logic")
            
            # Check for active time-off session or detect new time-off intent
            if self.timeoff_service and self.session_manager:
                try:
                    debug_log(f"Checking time-off flow for message: {message[:50]}...", "bot_logic")
                    # If no active session for provided thread_id, consider rebinding ONLY
                    # when the user message looks like a continuation (dates/yes/no/1/2/3 etc.).
                    def _looks_like_timeoff_continuation(txt: str) -> bool:
                        try:
                            s = (txt or '').strip().lower()
                            if not s:
                                return False
                            # clear intents should NOT trigger rebind
                            starters = ['time off', 'leave', 'annual leave', 'sick leave', 'request time off']
                            if any(k in s for k in starters):
                                return False
                            # continuation tokens
                            if s in {'yes','y','no','n','submit','confirm','cancel','stop','exit','quit'}:
                                return True
                            if any(tok in s for tok in ['hour_from=','hour_to=']):
                                return True
                            if any(k in s for k in ['annual','sick','custom hours']):
                                return True
                            # numeric choice (1,2,3)
                            if s.isdigit() and len(s) <= 2:
                                return True
                            import re as _re
                            if _re.search(r"\d{1,2}[/-]\d{1,2}(?:[/-]\d{2,4})?", s):
                                return True
                            if ' to ' in s or ' until ' in s or ' till ' in s or '-' in s:
                                # likely a range
                                return True
                            return False
                        except Exception:
                            return False

                    try:
                        active_for_thread = self.session_manager.get_active_session(thread_id) if thread_id else None
                    except Exception:
                        active_for_thread = None
                    if not active_for_thread and _looks_like_timeoff_continuation(message):
                        try:
                            active_list = self.session_manager.find_active_timeoff_sessions()
                            if isinstance(active_list, list) and len(active_list) == 1:
                                rebound_thread_id, rebound_session = active_list[0]
                                thread_id = thread_id or rebound_thread_id
                                debug_log(f"Rebinding to active time-off session thread: {thread_id}", "bot_logic")
                        except Exception:
                            pass

                    timeoff_response = self._handle_timeoff_flow(message, thread_id, employee_data)
                    if timeoff_response:
                        debug_log(f"Time-off flow handled, returning response", "bot_logic")
                        return timeoff_response
                    debug_log(f"Time-off flow did not handle message, continuing to normal flow", "bot_logic")
                except Exception as timeoff_error:
                    debug_log(f"Error in time-off flow: {timeoff_error}", "general")
                    import traceback
                    traceback.print_exc()
                    
                    # Clear any sessions that might be stuck
                    if thread_id:
                        self.session_manager.clear_session(thread_id)
                        
                    # Return a helpful error message instead of hanging
                    return {
                        'message': "I tried to help with your time-off request but encountered a technical issue. Please try again or contact HR for assistance. Error details have been logged for troubleshooting.",
                        'thread_id': thread_id,
                        'source': self.model,
                        'confidence_score': 1.0,
                        'model_used': self.model
                    }
            
            # Check for active reimbursement session or detect new reimbursement intent
            if self.reimbursement_service and self.session_manager:
                try:
                    debug_log(f"Checking reimbursement flow for message: {message[:50]}...", "bot_logic")
                    reimbursement_response = self.reimbursement_service.handle_flow(message, thread_id, employee_data)
                    if reimbursement_response:
                        debug_log(f"Reimbursement flow handled, returning response", "bot_logic")
                        return reimbursement_response
                    debug_log(f"Reimbursement flow did not handle message, continuing to normal flow", "bot_logic")
                except Exception as reimbursement_error:
                    debug_log(f"Error in reimbursement flow: {reimbursement_error}", "general")
                    import traceback
                    traceback.print_exc()
                    
                    # Clear any sessions that might be stuck
                    if thread_id:
                        self.session_manager.clear_session(thread_id)
                        
                    # Return a helpful error message instead of hanging
                    return {
                        'message': "I tried to help with your reimbursement request but encountered a technical issue. Please try again or contact HR for assistance. Error details have been logged for troubleshooting.",
                        'thread_id': thread_id,
                        'source': self.model,
                        'confidence_score': 1.0,
                        'model_used': self.model
                    }
            
            # Load conversation history and apply cap/condense strategy
            full_history = self._load_conversation_history(thread_id)
            max_history = getattr(Config, 'MAX_HISTORY_MESSAGES', 20)
            context_limit = getattr(Config, 'HISTORY_CONTEXT_LIMIT', 150000)
            include_condensed = getattr(Config, 'INCLUDE_CONDENSED_HISTORY', True)

            # Split into older and recent
            recent_history = full_history[-max_history:] if full_history else []
            older_history = full_history[:-max_history] if full_history and len(full_history) > max_history else []

            condensed_context = ""
            if older_history:
                # Condense older messages into a compact context paragraph
                try:
                    pairs = []
                    for h in older_history:
                        role = h.get('role')
                        content = h.get('content', '')
                        if role in ('user', 'assistant') and content:
                            tag = 'U' if role == 'user' else 'A'
                            pairs.append(f"[{tag}] {content}")
                    condensed_context = "\n".join(pairs)
                    if len(condensed_context) > context_limit:
                        condensed_context = condensed_context[:context_limit] + "\n..."
                except Exception:
                    condensed_context = ""
            
            # Prepare system message with Nasma rules and facts memory
            system_message = """You are Nasma, a helpful, precise assistant for PrezLab. 

When providing information about company policies, procedures, or guidelines, always give comprehensive, detailed explanations. Break down complex topics into clear sections and provide specific details, examples, and step-by-step processes when applicable.

For policy-related questions, include:
- Complete policy details and requirements
- Specific procedures and processes
- Applicable rates, calculations, or formulas
- Country-specific variations (Jordan, UAE, KSA)
- Approval processes and workflows
- Examples and scenarios when helpful

Be thorough and informative while maintaining clarity and accuracy."""

            # Facts memory (from Odoo + chat-derived preferences placeholder)
            facts_lines = []
            if employee_data:
                facts_lines.append(f"Name: {employee_data.get('name', 'Unknown')}")
                facts_lines.append(f"Job Title: {employee_data.get('job_title', 'Unknown')}")
                dept = employee_data.get('department_id_details', {})
                facts_lines.append(f"Department: {dept.get('name') if isinstance(dept, dict) else 'Unknown'}")
                manager = employee_data.get('parent_id_details', {})
                facts_lines.append(f"Manager: {manager.get('name') if isinstance(manager, dict) else 'Unknown'}")
                facts_lines.append(f"Time Zone: {employee_data.get('tz', 'Unknown')}")
                company = employee_data.get('company_id_details', {})
                facts_lines.append(f"Company: {company.get('name') if isinstance(company, dict) else 'Prezlab'}")

            facts_block = "\n".join(["Facts:"] + facts_lines) if facts_lines else ""

            # Rolling conversation summary
            rolling_summary = self._load_summary(thread_id)
            
            # Build messages array per requested layout
            messages = [{"role": "system", "content": system_message}]

            # Simple KB: read and inject as a single system message
            try:
                if False:  # Temporarily disable KB to fix hanging issue
                    kb_dir = os.path.normpath(os.path.join(os.path.dirname(__file__), '..', 'knowledge_base'))
                    max_chars = int(getattr(Config, 'KB_MAX_CHARS', 2000))
                    debug_log(f"Knowledge base enabled - checking directory: {kb_dir}", "knowledge_base")
                    debug_log(f"Max characters limit: {max_chars}", "knowledge_base")
                    
                    if os.path.isdir(kb_dir):
                        kb_texts = []
                        kb_used_files = []
                        all_files = os.listdir(kb_dir)
                        debug_log(f"Found {len(all_files)} files in knowledge base: {all_files}", "knowledge_base")
                        
                        for fname in all_files:
                            lower = fname.lower()
                            path = os.path.join(kb_dir, fname)
                            try:
                                if lower.endswith(('.md', '.txt')):
                                    debug_log(f"Processing text file: {fname}", "knowledge_base")
                                    with open(path, 'r', encoding='utf-8') as f:
                                        content = f.read()
                                        kb_texts.append(f"# {fname}\n" + content)
                                        kb_used_files.append(fname)
                                        debug_log(f"Loaded {fname}: {len(content)} characters", "knowledge_base")
                                elif lower.endswith('.docx'):
                                    try:
                                        debug_log(f"Processing DOCX file: {fname}", "knowledge_base")
                                        from docx import Document
                                        doc = Document(path)
                                        parts = [p.text for p in doc.paragraphs if p.text]
                                        content = "\n".join(parts)
                                        kb_texts.append(f"# {fname}\n" + content)
                                        kb_used_files.append(fname)
                                        debug_log(f"Loaded {fname}: {len(content)} characters", "knowledge_base")
                                    except Exception as de:
                                        debug_log(f"DOCX parse error for {fname}: {de}", "knowledge_base")
                                elif lower.endswith('.pptx'):
                                    try:
                                        debug_log(f"Processing PPTX file: {fname}", "knowledge_base")
                                        from pptx import Presentation
                                        prs = Presentation(path)
                                        slides_text = []
                                        for s in prs.slides:
                                            buf = []
                                            for shp in s.shapes:
                                                if hasattr(shp, 'text') and shp.text:
                                                    buf.append(shp.text)
                                            if buf:
                                                slides_text.append("\n".join(buf))
                                        content = "\n\n".join(slides_text)
                                        kb_texts.append(f"# {fname}\n" + content)
                                        kb_used_files.append(fname)
                                        debug_log(f"Loaded {fname}: {len(content)} characters", "knowledge_base")
                                    except Exception as pe:
                                        debug_log(f"PPTX parse error for {fname}: {pe}", "knowledge_base")
                            except Exception as fe:
                                debug_log(f"KB read error for {fname}: {fe}", "knowledge_base")
                        
                        if kb_texts:
                            # Prioritize important files when truncating
                            priority_files = ['overtime', 'policy', 'leave', 'holiday']
                            priority_texts = []
                            regular_texts = []
                            
                            for text in kb_texts:
                                is_priority = any(priority in text.lower() for priority in priority_files)
                                if is_priority:
                                    priority_texts.append(text)
                                else:
                                    regular_texts.append(text)
                            
                            # Combine with priority files first
                            ordered_texts = priority_texts + regular_texts
                            kb_blob = "\n\n".join(ordered_texts)
                            debug_log(f"Total knowledge base content: {len(kb_blob)} characters", "knowledge_base")
                            debug_log(f"Priority files found: {len(priority_texts)}", "knowledge_base")
                            
                            if len(kb_blob) > max_chars:
                                # Try to keep priority content
                                if priority_texts:
                                    priority_blob = "\n\n".join(priority_texts)
                                    remaining_chars = max_chars - len(priority_blob) - 100  # Reserve space
                                    if remaining_chars > 0:
                                        regular_blob = "\n\n".join(regular_texts)
                                        if len(regular_blob) > remaining_chars:
                                            regular_blob = regular_blob[:remaining_chars] + "\n..."
                                        kb_blob = priority_blob + "\n\n" + regular_blob
                                    else:
                                        kb_blob = priority_blob[:max_chars-3] + "..."
                                else:
                                    kb_blob = kb_blob[:max_chars] + "\n..."
                                
                                debug_log(f"KB content truncated to {len(kb_blob)} characters (prioritized)", "knowledge_base")
                                debug_log(f"KB files loaded (truncated): {kb_used_files}", "knowledge_base")
                            else:
                                debug_log(f"KB files loaded (full): {kb_used_files}", "knowledge_base")
                            
                            debug_log(f"Adding knowledge base to ChatGPT context with {len(kb_blob)} characters", "knowledge_base")
                            messages.append({"role": "system", "content": "KB:\n" + kb_blob})
                        else:
                            debug_log("No knowledge base content loaded", "knowledge_base")
                    else:
                        debug_log(f"KB directory not found: {kb_dir}", "knowledge_base")
                else:
                    debug_log("Knowledge base is disabled in configuration", "knowledge_base")
            except Exception as kb_e:
                debug_log(f"KB injection failed: {kb_e}", "knowledge_base")

            if facts_block:
                messages.append({"role": "system", "content": facts_block})

            if rolling_summary:
                messages.append({"role": "system", "content": "Conversation summary:\n" + rolling_summary})

            # Add only the recent capped history
            for msg in recent_history:
                messages.append({"role": msg["role"], "content": msg["content"]})
            
            # Add current user message
            messages.append({"role": "user", "content": message})
            
            debug_log(f"Sending {len(messages)} messages to {self.model}", "bot_logic")
            
            # Log what messages are being sent to ChatGPT
            for i, msg in enumerate(messages):
                role = msg.get('role', 'unknown')
                content_preview = msg.get('content', '')[:200] + "..." if len(msg.get('content', '')) > 200 else msg.get('content', '')
                debug_log(f"Message {i+1} ({role}): {content_preview}", "knowledge_base")
            
            # Build request args depending on model capability
            def _build_chat_args(model_name: str, temperature: float = 1):
                base_args = {
                    'model': model_name,
                    'messages': messages,
                    'temperature': temperature,
                    'top_p': 1.0,
                    'frequency_penalty': 0.0,
                    'presence_penalty': 0.0,
                    'response_format': {'type': 'text'}
                }
                # gpt-5 uses max_completion_tokens instead of max_tokens per error message
                if model_name.startswith('gpt-5'):
                    base_args['max_completion_tokens'] = 2000
                else:
                    base_args['max_tokens'] = 2000
                return base_args

            # Make API call to the configured model with optional fallback
            try:
                debug_log(f"Making API call to {self.model}...", "bot_logic")
                client = self._get_client()
                if client is None:
                    return {
                        'message': "I'm experiencing technical difficulties. Please try again later.",
                        'thread_id': thread_id,
                        'source': self.model,
                        'confidence_score': 0.0,
                        'model_used': self.model
                    }
                response = client.chat.completions.create(**_build_chat_args(self.model), timeout=30)
                debug_log(f"API call completed successfully", "bot_logic")
            except Exception as primary_error:
                if self.fallback_model:
                    debug_log(f"Primary model {self.model} failed: {primary_error}. Falling back to {self.fallback_model}", "bot_logic")
                    client = self._get_client()
                    if client is None:
                        return {
                            'message': "I'm experiencing technical difficulties. Please try again later.",
                            'thread_id': thread_id,
                            'source': self.model,
                            'confidence_score': 0.0,
                            'model_used': self.model
                        }
                    response = client.chat.completions.create(**_build_chat_args(self.fallback_model))
                    # Update reported model
                    self.model = self.fallback_model
                else:
                    raise
            
            # Log finish_reason and usage for diagnostics
            try:
                finish_reason = getattr(response.choices[0], 'finish_reason', None)
                usage = getattr(response, 'usage', None)
                debug_log(f"finish_reason={finish_reason}, usage={usage}", "bot_logic")
            except Exception:
                pass

            # Extract response (guard against empty content)
            response_text = getattr(response.choices[0].message, 'content', None)
            if response_text is None or (isinstance(response_text, str) and response_text.strip() == ""):
                # Auto-retry once with safer prompt and lower temperature
                debug_log("Empty response, retrying once with safer prompt and lower temperature", "bot_logic")
                retry_messages = list(messages)
                retry_messages.insert(0, {"role": "system", "content": "Reply in plain text only. If any part seems unsafe or unclear, provide a brief safe explanation instead of returning nothing."})
                try:
                    client = self._get_client()
                    if client is None:
                        return {
                            'message': "I'm experiencing technical difficulties. Please try again later.",
                            'thread_id': thread_id,
                            'source': self.model,
                            'confidence_score': 0.0,
                            'model_used': self.model
                        }
                    retry_response = client.chat.completions.create(**_build_chat_args(self.model, temperature=0.3) | {'messages': retry_messages})
                    try:
                        finish_reason_retry = getattr(retry_response.choices[0], 'finish_reason', None)
                        usage_retry = getattr(retry_response, 'usage', None)
                        debug_log(f"retry finish_reason={finish_reason_retry}, usage={usage_retry}", "bot_logic")
                    except Exception:
                        pass
                    response_text = getattr(retry_response.choices[0].message, 'content', None)
                except Exception as retry_error:
                    debug_log(f"Retry failed: {retry_error}", "general")
                
                # Final fallback text if still empty
                if response_text is None or (isinstance(response_text, str) and response_text.strip() == ""):
                    response_text = "Sorry, I couldn't generate a response just now. Please try rephrasing or asking again."
            
            # Update conversation history (store full turns; condensing happens on send)
            updated_history = list(full_history) if full_history else []
            updated_history.append({"role": "user", "content": message})
            updated_history.append({"role": "assistant", "content": response_text})

            # Update rolling summary (simple heuristic: keep last 1000 chars combining recent messages)
            try:
                summary_source = rolling_summary + "\n" + "\n".join(
                    [f"U: {m['content']}" for m in recent_history if m.get('role') == 'user'] +
                    [f"A: {m['content']}" for m in recent_history if m.get('role') == 'assistant'] +
                    [f"A: {response_text}"]
                )
                new_summary = summary_source[-1000:]
                self._save_summary(thread_id, new_summary)
            except Exception as e:
                debug_log(f"Failed to update rolling summary: {e}", "general")
            
            # Save updated conversation history
            self._save_conversation_history(thread_id, updated_history)
            
            debug_log(f"{self.model} response: '{response_text[:100]}...'", "bot_logic")
                    
            return {
                'message': response_text,
                'thread_id': thread_id,
                'source': self.model,
                'confidence_score': 1.0,
                'model_used': self.model
            }
                
        except Exception as e:
            error_msg = f"Error communicating with {self.model}: {str(e)}"
            debug_log(f"{error_msg}", "general")
            return {
                'message': error_msg,
                'thread_id': thread_id,
                'source': self.model,
                'confidence_score': 0.0,
                'error': True
            }
    
    def clear_conversation_history(self, thread_id):
        """Clear conversation history for a specific thread"""
        try:
            with self.lock:
                storage_file = self._get_storage_file(thread_id)
                if os.path.exists(storage_file):
                    os.remove(storage_file)
                    debug_log(f"Cleared conversation history for thread: {thread_id}", "bot_logic")
                    return True
                return False
        except Exception as e:
            debug_log(f"Error clearing conversation history: {e}", "general")
            return False
    
    def get_conversation_history(self, thread_id):
        """Get conversation history for a specific thread"""
        return self._load_conversation_history(thread_id)
    
    def get_model_info(self):
        """Get information about the current model being used"""
        return {
            'model': self.model,
            'description': f'{self.model} - Latest available OpenAI model with enhanced capabilities',
            'features': ['Text generation', 'Context understanding', 'Employee personalization', 'Conversation history']
        }
    
    def _handle_timeoff_flow(self, message: str, thread_id: str, employee_data: dict) -> dict:
        """Handle time-off request flow with session management"""
        try:
            debug_log(f"Starting time-off flow check...", "bot_logic")

            # Enforce single active flow per thread: block if another flow is active
            try:
                active_any = self.session_manager.get_active_session(thread_id) if thread_id else None
                if active_any and active_any.get('type') not in (None, 'timeoff') and active_any.get('state') in ['started', 'active']:
                    other = active_any.get('type', 'another')
                    return {
                        'message': f"You're currently in an active {other} request. Please complete it or type 'cancel' to end it before starting time off.",
                        'thread_id': thread_id,
                        'source': self.model,
                        'confidence_score': 1.0,
                        'session_handled': True
                    }
            except Exception:
                pass

            # Check for active session using thread_id if provided
            active_session = None
            if thread_id:
                active_session = self.session_manager.get_session(thread_id)
                debug_log(f"Active session check for thread_id {thread_id}: {active_session is not None}", "bot_logic")

            # Detect new time-off intent first
            is_timeoff, confidence, extracted_data = self.timeoff_service.detect_timeoff_intent(message)
            debug_log(f"Time-off detection complete: is_timeoff={is_timeoff}, confidence={confidence}", "bot_logic")

            # If this is a new time-off request (high confidence)
            if is_timeoff and confidence >= 0.7:
                # If there's already an active time-off session, continue it instead of restarting
                if active_session and active_session.get('type') == 'timeoff' and active_session.get('state') in ['started', 'active']:
                    debug_log(f"Active time-off session exists; continuing instead of restarting", "bot_logic")
                    return self._continue_timeoff_session(message, thread_id, active_session, employee_data)

                debug_log(f"High confidence time-off intent detected ({confidence:.2f}), starting fresh session...", "bot_logic")
                # Clear any orphaned active sessions for this user to start fresh
                if active_session:
                    self.session_manager.clear_session(thread_id)
                # Also clear any orphaned active sessions to prevent conflicts
                active_sessions = self.session_manager.find_active_timeoff_sessions()
                for orphaned_thread_id, orphaned_session in active_sessions:
                    if orphaned_session.get('state') in ['started', 'active']:
                        debug_log(f"Clearing orphaned session: {orphaned_thread_id}", "bot_logic")
                        self.session_manager.clear_session(orphaned_thread_id)
                # Validate that we have employee data before starting session
                if not employee_data or not isinstance(employee_data, dict) or not employee_data.get('id'):
                    debug_log(f"Invalid employee data for time-off request: {employee_data}", "bot_logic")
                    return {
                        'message': "I'd like to help with your time-off request, but I need to verify your employee information first. Please try logging out and logging back in, or contact HR for assistance.",
                        'thread_id': thread_id,
                        'source': self.model,
                        'confidence_score': 1.0,
                        'model_used': self.model
                    }

                debug_log(f"Time-off intent detected with confidence {confidence:.2f}, starting session...", "bot_logic")
                timeoff_result = self._start_timeoff_session(message, thread_id, extracted_data, employee_data)
                debug_log(f"Time-off session start result: {timeoff_result is not None}", "bot_logic")
                if timeoff_result:
                    debug_log(f"Time-off session returned result with message: {timeoff_result.get('message', 'NO MESSAGE')[:50]}...", "bot_logic")
                else:
                    debug_log(f"Time-off session returned None - this is the problem!", "bot_logic")
                return timeoff_result

            # If we have an active session, continue it
            if active_session and active_session.get('type') == 'timeoff':
                # Validate session state - accept both 'started' and 'active' as valid
                session_state = active_session.get('state', 'unknown')
                if session_state in ['completed', 'cancelled']:
                    debug_log(f"Session is {session_state}, clearing and returning None", "bot_logic")
                    self.session_manager.clear_session(thread_id)
                    return None
                elif session_state in ['started', 'active']:
                    debug_log(f"Found active time-off session in step {active_session.get('step', 'unknown')}, continuing...", "bot_logic")
                    return self._continue_timeoff_session(message, thread_id, active_session, employee_data)

            # Lower confidence time-off detection (0.3 - 0.7)
            if is_timeoff and confidence >= 0.3:
                # Validate that we have employee data before starting session
                if not employee_data or not isinstance(employee_data, dict) or not employee_data.get('id'):
                    debug_log(f"Invalid employee data for time-off request: {employee_data}", "bot_logic")
                    return {
                        'message': "I'd like to help with your time-off request, but I need to verify your employee information first. Please try logging out and logging back in, or contact HR for assistance.",
                        'thread_id': thread_id,
                        'source': self.model,
                        'confidence_score': 1.0,
                        'model_used': self.model
                    }

                debug_log(f"Time-off intent detected with confidence {confidence:.2f}, starting session...", "bot_logic")
                timeoff_result = self._start_timeoff_session(message, thread_id, extracted_data, employee_data)
                debug_log(f"Time-off session start result: {timeoff_result is not None}", "bot_logic")
                if timeoff_result:
                    debug_log(f"Time-off session returned result with message: {timeoff_result.get('message', 'NO MESSAGE')[:50]}...", "bot_logic")
                else:
                    debug_log(f"Time-off session returned None - this is the problem!", "bot_logic")
                return timeoff_result

            debug_log(f"No time-off intent detected (confidence: {confidence:.2f}), returning None", "bot_logic")
            return None

        except Exception as e:
            debug_log(f"Error in time-off flow: {e}", "general")
            import traceback
            traceback.print_exc()
            # Clear any existing session that might be causing issues
            if thread_id:
                self.session_manager.clear_session(thread_id)
            return None
    
    def _start_timeoff_session(self, message: str, thread_id: str, extracted_data: dict, employee_data: dict) -> dict:
        """Start a new time-off request session"""
        try:
            # Validate thread_id
            if not thread_id:
                debug_log("Warning: thread_id is None or empty, generating fallback ID", "bot_logic")
                import time
                thread_id = f"timeoff_{int(time.time())}"
                debug_log(f"Generated fallback thread_id: {thread_id}", "bot_logic")
            
            debug_log(f"Starting time-off session with thread_id: {thread_id}", "bot_logic")
            
            # Clear any existing session first to prevent conflicts
            self.session_manager.clear_session(thread_id)
            
            # Start session
            session_data = {
                'extracted_data': extracted_data,
                'employee_data': employee_data
            }
            session = self.session_manager.start_session(thread_id, 'timeoff', session_data)
            
            # Get available leave types from Odoo
            debug_log(f"Fetching leave types for time-off session...", "bot_logic")
            success, leave_types = self.timeoff_service.get_leave_types()
            debug_log(f"Leave types fetch result: success={success}, data={leave_types}", "bot_logic")
            
            if not success:
                # Handle Odoo connection issues gracefully
                debug_log(f"Failed to fetch leave types from Odoo: {leave_types}", "bot_logic")
                self.session_manager.clear_session(thread_id)
                return {
                    'message': "I'd like to help with your time-off request, but I'm unable to connect to the HR system right now. Please try again later or contact HR directly for assistance.",
                    'thread_id': thread_id,
                    'source': self.model,
                    'confidence_score': 1.0,
                    'model_used': self.model
                }
            
            # Validate and clean leave types data
            if not leave_types or not isinstance(leave_types, list):
                debug_log(f"Invalid leave types data received: {leave_types}", "bot_logic")
                self.session_manager.clear_session(thread_id)
                return {
                    'message': "I'm having trouble accessing the available leave types. Please contact HR for assistance with your time-off request.",
                    'thread_id': thread_id,
                    'source': self.model,
                    'confidence_score': 1.0,
                    'model_used': self.model
                }
            
            # Clean corrupted leave types data
            debug_log(f"Cleaning {len(leave_types)} leave types for corruption...", "bot_logic")
            clean_leave_types = []
            seen_names = set()
            
            for i, lt in enumerate(leave_types):
                try:
                    # Skip if not a proper dictionary
                    if not isinstance(lt, dict):
                        debug_log(f"Skipping non-dict leave type at index {i}: {type(lt)}", "bot_logic")
                        continue
                    
                    # Must have required fields
                    if 'id' not in lt or 'name' not in lt:
                        debug_log(f"Skipping incomplete leave type at index {i}: {lt}", "bot_logic")
                        continue
                    
                    # Skip duplicates by name
                    name = lt.get('name', '').strip()
                    if not name or name in seen_names:
                        debug_log(f"Skipping duplicate/empty name at index {i}: {name}", "bot_logic")
                        continue
                    
                    # Create clean entry with only essential fields
                    clean_entry = {
                        'id': lt.get('id'),
                        'name': name,
                        'active': lt.get('active', True)
                    }
                    
                    clean_leave_types.append(clean_entry)
                    seen_names.add(name)
                    
                except Exception as clean_error:
                    debug_log(f"Error cleaning leave type {i}: {clean_error}", "bot_logic")
                    continue
            
            debug_log(f"Cleaned leave types: {len(clean_leave_types)} valid entries", "bot_logic")
            
            if not clean_leave_types:
                debug_log(f"No valid leave types after cleaning", "bot_logic")
                self.session_manager.clear_session(thread_id)
                return {
                    'message': "The leave types data from HR system appears to be corrupted. Please contact HR directly for assistance with your time-off request.",
                    'thread_id': thread_id,
                    'source': self.model,
                    'confidence_score': 1.0,
                    'model_used': self.model
                }
            
            # Use cleaned data
            leave_types = clean_leave_types

            # Inject Half Days option by replacing Unpaid Leave when applicable
            try:
                if self.halfday_service:
                    leave_types = self.halfday_service.replace_unpaid_with_halfdays(leave_types)
                    debug_log(f"HalfDay: Post-replacement leave types count: {len(leave_types)}", "bot_logic")
            except Exception as hd_e:
                debug_log(f"HalfDay replacement failed: {hd_e}", "general")
            
            # Update session with leave types - with error handling
            try:
                debug_log(f"Updating session with {len(leave_types)} leave types", "bot_logic")
                debug_log(f"About to call session_manager.update_session...", "bot_logic")
                self.session_manager.update_session(thread_id, {'leave_types': leave_types})
                debug_log(f"Session updated successfully", "bot_logic")
            except Exception as session_error:
                debug_log(f"Error updating session: {session_error}", "general")
                import traceback
                traceback.print_exc()
                self.session_manager.clear_session(thread_id)
                return {
                    'message': "I'm having trouble managing your time-off request session. Please try again.",
                    'thread_id': thread_id,
                    'source': self.model,
                    'confidence_score': 1.0,
                    'model_used': self.model
                }
            
            # Check if leave type was already extracted
            try:
                debug_log(f"Checking for pre-extracted leave type...", "bot_logic")
                if 'leave_type' in extracted_data:
                    debug_log(f"Found pre-extracted leave type: {extracted_data['leave_type']}", "bot_logic")
                    # Try to map to actual leave type
                    extracted_type = extracted_data['leave_type']
                    matched_type = None
                    for lt in leave_types:
                        if extracted_type.lower() in lt.get('name', '').lower():
                            matched_type = lt
                            break
                    
                    if matched_type:
                        debug_log(f"Matched leave type: {matched_type['name']}", "bot_logic")
                        self.session_manager.update_session(thread_id, {'selected_leave_type': matched_type})
                        self.session_manager.advance_session_step(thread_id, {'leave_type_confirmed': True})

                        # If both dates are already present, move straight to confirmation
                        if 'start_date' in extracted_data and 'end_date' in extracted_data:
                            start_date = self.timeoff_service.parse_date_input(extracted_data['start_date']) or extracted_data['start_date']
                            end_date = self.timeoff_service.parse_date_input(extracted_data['end_date']) or extracted_data['end_date']
                            self.session_manager.update_session(thread_id, {'start_date': start_date, 'end_date': end_date})
                            self.session_manager.advance_session_step(thread_id)

                            def dd_slash_mm_yyyy(d: str) -> str:
                                try:
                                    return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
                                except Exception:
                                    return d
                            response_text = f"Perfect! Here's your time-off request summary:\n\n"
                            response_text += f"📋 **Leave Type:** {matched_type.get('name', 'Unknown')}\n"
                            response_text += f"📅 **Start Date:** {dd_slash_mm_yyyy(start_date)}\n"
                            response_text += f"📅 **End Date:** {dd_slash_mm_yyyy(end_date)}\n"
                            response_text += f"👤 **Employee:** {employee_data.get('name', 'Unknown')}\n\n"
                            response_text += "Do you want to submit this request? Reply with 'yes' to confirm or 'no' to cancel."
                            return self._create_response(response_text, thread_id)

                        # Otherwise, ask for both dates in a single message
                        response_text = (
                            f"I'll help you request {matched_type['name']}. \n\n"
                            "You can pick dates from the calendar below or type them. Examples:\n"
                            "- 23/9 to 24/9\n"
                            "- 23/09/2025 to 24/09/2025\n"
                            "- 23-9-2025 till 24-9-2025\n"
                            "- 23rd of September till the 24th\n"
                            "- next Monday to Wednesday\n\n"
                            "Defaults: I assume DD/MM, current month and year unless you specify otherwise."
                        )
                        return self._create_response_with_datepicker(response_text, thread_id)
                else:
                    debug_log(f"No pre-extracted leave type found", "bot_logic")
            except Exception as extract_error:
                debug_log(f"Error processing extracted data: {extract_error}", "general")
                # Continue to present options
            
            # Present the 3 main leave type options with buttons
            try:
                debug_log(f"Presenting main leave type options to user...", "bot_logic")

                # Filter to only show the 3 main leave types
                main_leave_types = []
                # Replace 'Unpaid Leave' with 'Custom Hours'
                type_names = ['Annual Leave', 'Sick Leave', 'Custom Hours']

                for type_name in type_names:
                    for lt in leave_types:
                        if lt.get('name') == type_name:
                            main_leave_types.append(lt)
                            break

                debug_log(f"Found {len(main_leave_types)} main leave types", "bot_logic")

                # Store the main leave types in session
                self.session_manager.update_session(thread_id, {'main_leave_types': main_leave_types})

                response_text = "I'll help you request time off! Please select the type of leave you need:"

                # Create response with buttons
                result = self._create_response_with_buttons(
                    response_text,
                    thread_id,
                    main_leave_types
                )
                debug_log(f"Response with buttons created successfully: {result is not None}", "bot_logic")
                return result
                
            except Exception as response_error:
                debug_log(f"Error creating leave type options response: {response_error}", "general")
                import traceback
                traceback.print_exc()
                
                # Fallback response
                self.session_manager.clear_session(thread_id)
                return {
                    'message': "I can help you request time off, but I'm having trouble loading the available leave types right now. Please contact HR for assistance.",
                    'thread_id': thread_id,
                    'source': self.model,
                    'confidence_score': 1.0,
                    'model_used': self.model
                }
            
        except Exception as e:
            debug_log(f"Error starting time-off session: {e}", "general")
            # Clear session on error to prevent stuck state
            self.session_manager.clear_session(thread_id)
            return {
                'message': "I'd like to help with your time-off request, but I'm experiencing some technical difficulties. Please try again later.",
                'thread_id': thread_id,
                'source': self.model,
                'confidence_score': 1.0,
                'model_used': self.model
            }
    
    def _continue_timeoff_session(self, message: str, thread_id: str, session: dict, employee_data: dict) -> dict:
        """Continue an active time-off session"""
        try:
            # If session is already completed/cancelled, clear and do not continue
            try:
                state = session.get('state')
                if state in ['completed', 'cancelled']:
                    self.session_manager.clear_session(thread_id)
                    return None
            except Exception:
                pass
            # Check for exit commands
            message_lower = message.lower().strip()
            if message_lower in ['cancel', 'exit', 'stop', 'quit', 'nevermind', 'no thanks', 'end', 'abort', 'undo', 'no', 'n']:
                debug_log(f"User said '{message_lower}', cancelling time-off session", "bot_logic")
                try:
                    self.session_manager.cancel_session(thread_id, f"User requested to exit time-off flow with: {message_lower}")
                finally:
                    # Clear any persisted state to avoid sticky sessions
                    self.session_manager.clear_session(thread_id)
                return {
                    'message': 'request cancelled, can i help you with anything else',
                    'thread_id': thread_id,
                    'source': self.model,
                    'confidence_score': 1.0,
                    'model_used': self.model
                }

            step = session.get('step', 1)
            debug_log(f"Continuing session at step {step} for thread {thread_id}", "bot_logic")

            if step == 1:  # Waiting for leave type selection
                return self._handle_leave_type_selection(message, thread_id, session, employee_data)
            elif step == 2:  # Waiting for date range (start and end in one message)
                # Safety: if no leave type was selected (e.g., residue from a previous flow), reset to step 1
                try:
                    sd = session.get('data', {}) if isinstance(session, dict) else {}
                    selected_type = sd.get('selected_leave_type') or session.get('selected_leave_type')
                except Exception:
                    selected_type = None
                if not selected_type:
                    debug_log("Step=2 but no selected_leave_type found; resetting to step 1", "bot_logic")
                    self.session_manager.update_session(thread_id, {'step': 1})
                    return self._handle_leave_type_selection(message, thread_id, session, employee_data)
                return self._handle_date_range_input(message, thread_id, session, employee_data)
            elif step == 3:  # Waiting for confirmation
                return self._handle_confirmation(message, thread_id, session, employee_data)
            else:
                debug_log(f"Session in unknown state (step {step}), restarting", "bot_logic")
                # Session in unknown state, restart
                self.session_manager.clear_session(thread_id)
                return None

        except Exception as e:
            debug_log(f"Error continuing time-off session: {e}", "general")
            import traceback
            traceback.print_exc()
            # Clear the broken session
            self.session_manager.clear_session(thread_id)
            return None
    
    def _handle_leave_type_selection(self, message: str, thread_id: str, session: dict, employee_data: dict) -> dict:
        """Handle leave type selection step"""
        session_data = session.get('data', {})
        # Check for main leave types first (the 3 button options)
        leave_types = session_data.get('main_leave_types', [])
        if not leave_types:
            # Fallback to full leave types
            leave_types = session_data.get('leave_types', [])
        if not leave_types:
            # Also check at session root level
            leave_types = session.get('leave_types', []) or session.get('main_leave_types', [])

        debug_log(f"Found {len(leave_types)} leave types in session", "bot_logic")
        
        # Validate that we still have leave types
        if not leave_types:
            debug_log("No leave types found in session data, cancelling session", "bot_logic")
            self.session_manager.clear_session(thread_id)
            return {
                'message': "I'm sorry, but I'm having trouble accessing the leave types. Please try requesting time off again or contact HR for assistance.",
                'thread_id': thread_id,
                'source': self.model,
                'confidence_score': 1.0,
                'model_used': self.model
            }
        
        message_clean = message.strip()
        
        # Check if user provided a number
        try:
            choice_num = int(message_clean)
            if 1 <= choice_num <= len(leave_types):
                selected_type = leave_types[choice_num - 1]
                self.session_manager.update_session(thread_id, {'selected_leave_type': selected_type})
                self.session_manager.advance_session_step(thread_id)
                
                response_text = f"Great! You've selected {selected_type['name']}. \n\nYou can pick dates from the calendar below or type them. Examples:\n- 23/9 to 24/9\n- 23/09/2025 to 24/09/2025\n- 23-9-2025 till 24-9-2025\n- 23rd of September till the 24th\n- next Monday to Wednesday\n\nDefaults: I assume DD/MM, current month and year unless you specify otherwise."
                return self._create_response_with_datepicker(response_text, thread_id)
        except ValueError:
            pass
        
        # Try to match by name (improved matching)
        message_lower = message_clean.lower()
        best_match = None
        best_score = 0
        
        for leave_type in leave_types:
            leave_name = leave_type.get('name', '').lower()
            # Check for exact word matches
            name_words = leave_name.split()
            message_words = message_lower.split()
            
            matches = 0
            for word in name_words:
                if word in message_words:
                    matches += 1
            
            # Also check for partial name inclusion
            if leave_name in message_lower or message_lower in leave_name:
                matches += len(name_words)
            
            if matches > best_score:
                best_score = matches
                best_match = leave_type
        
        # If we found a good match, use it
        if best_match and best_score > 0:
            self.session_manager.update_session(thread_id, {'selected_leave_type': best_match})
            self.session_manager.advance_session_step(thread_id)
            
            # If Half Days, show single-date picker; else show range picker
            try:
                is_halfday = False
                if self.halfday_service and isinstance(best_match, dict):
                    is_halfday = self.halfday_service.is_halfday(best_match) or (best_match.get('name') == getattr(self.halfday_service, 'HALF_DAY_NAME', 'Custom Hours'))
            except Exception:
                is_halfday = False

            if is_halfday:
                response_text = (
                    f"Perfect! You've selected {best_match['name']}. \n\n"
                    "Pick your date from the calendar below or type it (single day only). Examples:\n"
                    "- 23/9\n"
                    "- 23/09/2025\n"
                    "- next Monday\n\n"
                    "Defaults: I assume DD/MM, current month and year unless you specify otherwise."
                )
                return self._create_response_with_datepicker_single(response_text, thread_id)
            else:
                response_text = (
                    f"Perfect! You've selected {best_match['name']}. \n\nYou can pick dates from the calendar below or type them. Examples:\n"
                    "- 23/9 to 24/9\n"
                    "- 23/09/2025 to 24/09/2025\n"
                    "- 23-9-2025 till 24-9-2025\n"
                    "- 23rd of September till the 24th\n"
                    "- next Monday to Wednesday\n\n"
                    "Defaults: I assume DD/MM, current month and year unless you specify otherwise."
                )
                return self._create_response_with_datepicker(response_text, thread_id)
        
        # No match found - before falling back, check if the user already provided dates.
        try:
            # If the message looks like a date range or a single date, treat it as step 2 input.
            dr = self.timeoff_service.parse_date_range(message)
            if dr:
                return self._handle_date_range_input(message, thread_id, session, employee_data)
            single = self.timeoff_service.parse_date_input(message)
            if single:
                # Treat as same-day range
                return self._handle_date_range_input(message, thread_id, session, employee_data)
        except Exception:
            pass

        # Still no match - provide helpful guidance
        leave_types_text = self.timeoff_service.format_leave_types_for_user(leave_types)
        response_text = f"I didn't quite catch that. Please select from the available options:\n\n{leave_types_text}\nYou can type the number (1, 2, 3, etc.) or the name of the leave type (like 'annual' or 'sick')."
        return self._create_response(response_text, thread_id)
    
    def _handle_start_date_input(self, message: str, thread_id: str, session: dict, employee_data: dict) -> dict:
        """Handle start date input step"""
        parsed_date = self.timeoff_service.parse_date_input(message)
        
        if parsed_date:
            self.session_manager.update_session(thread_id, {'start_date': parsed_date})
            self.session_manager.advance_session_step(thread_id)
            
            response_text = f"Got it! Start date: {parsed_date}\n\nWhat's your end date? Please use DD/MM/YYYY format."
            return self._create_response(response_text, thread_id)
        else:
            response_text = "I couldn't understand that date format. Please provide the start date in DD/MM/YYYY format (e.g., 25/12/2024)."
            return self._create_response(response_text, thread_id)
    
    def _handle_end_date_input(self, message: str, thread_id: str, session: dict, employee_data: dict) -> dict:
        """Handle end date input step"""
        parsed_date = self.timeoff_service.parse_date_input(message)
        
        if parsed_date:
            session_data = session.get('data', {})
            # Check both locations for start_date and selected_leave_type
            start_date = session_data.get('start_date') or session.get('start_date')
            selected_type = session_data.get('selected_leave_type') or session.get('selected_leave_type', {})
            # Robustly resolve employee data from argument or session
            resolved_employee = employee_data or session_data.get('employee_data') or session.get('employee_data') or {}
            
            debug_log(f"End date processing - start_date: {start_date}, selected_type: {selected_type.get('name', 'None') if selected_type else 'None'}", "bot_logic")
            
            # Validate end date is after start date
            if start_date and parsed_date < start_date:
                response_text = f"The end date ({parsed_date}) cannot be before the start date ({start_date}). Please provide a valid end date."
                return self._create_response(response_text, thread_id)
            
            self.session_manager.update_session(thread_id, {'end_date': parsed_date})
            self.session_manager.advance_session_step(thread_id)
            
            # Show confirmation (display dates as DD/MM/YYYY)
            def dd_mm_yyyy(d: str) -> str:
                try:
                    return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
                except Exception:
                    return d
            response_text = f"Perfect! Here's your time-off request summary:\n\n"
            response_text += f"📋 **Leave Type:** {selected_type.get('name', 'Unknown') if selected_type else 'Unknown'}\n"
            response_text += f"📅 **Start Date:** {dd_mm_yyyy(start_date) if start_date else 'Unknown'}\n"
            response_text += f"📅 **End Date:** {dd_mm_yyyy(parsed_date)}\n"
            response_text += f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}\n\n"
            response_text += "Do you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
            buttons = [
                {'text': 'Yes', 'value': 'yes', 'type': 'confirmation_choice'},
                {'text': 'No', 'value': 'no', 'type': 'confirmation_choice'}
            ]
            debug_log(f"Returning confirmation response with buttons: {response_text[:100]}...", "bot_logic")
            return self._create_response_with_choice_buttons(response_text, thread_id, buttons)
        else:
            response_text = "I couldn't understand that date format. Please provide the end date in DD/MM/YYYY format (e.g., 27/12/2024)."
            return self._create_response(response_text, thread_id)

    def _handle_date_range_input(self, message: str, thread_id: str, session: dict, employee_data: dict) -> dict:
        """Handle combined date range input (start and end in one message)"""
        # Detect if current flow is Half Days to enforce single-day constraint
        session_data_for_type = session.get('data', {})
        selected_type_for_validation = session_data_for_type.get('selected_leave_type') or session.get('selected_leave_type', {})
        is_halfday_flow = False
        try:
            if self.halfday_service and isinstance(selected_type_for_validation, dict):
                is_halfday_flow = self.halfday_service.is_halfday(selected_type_for_validation) or (
                    selected_type_for_validation.get('name') == getattr(self.halfday_service, 'HALF_DAY_NAME', 'Custom Hours')
                )
        except Exception:
            is_halfday_flow = False

        result = self.timeoff_service.parse_date_range(message)
        if result:
            start_date, end_date = result
            # If Half Day flow, only allow single day
            if is_halfday_flow and start_date != end_date:
                response_text = (
                    "For Half Days, you can only select one day. Please pick a single date."
                )
                return self._create_response_with_datepicker_single(response_text, thread_id)
            # Validate chronological order (already ensured) and store
            self.session_manager.update_session(thread_id, {'start_date': start_date, 'end_date': end_date})
            self.session_manager.advance_session_step(thread_id)

            session_data = session.get('data', {})
            selected_type = session_data.get('selected_leave_type') or session.get('selected_leave_type', {})
            resolved_employee = employee_data or session_data.get('employee_data') or session.get('employee_data') or {}

            # Display dates as DD/MM/YYYY
            def dd_mm_yyyy(d: str) -> str:
                try:
                    return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
                except Exception:
                    return d
            if is_halfday_flow:
                # After selecting the date for Half Days, ask for hour range
                hour_text = (
                    "Great, got your date. Please choose your hours (from/to)."
                )
                return self._create_response_with_hour_picker(hour_text, thread_id)
            else:
                response_text = f"Great, noted your dates. Here's your time-off request summary:\n\n"
                response_text += f"📋 **Leave Type:** {selected_type.get('name', 'Unknown') if selected_type else 'Unknown'}\n"
                response_text += f"📅 **Start Date:** {dd_mm_yyyy(start_date)}\n"
                response_text += f"📅 **End Date:** {dd_mm_yyyy(end_date)}\n"
                response_text += f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}\n\n"
                response_text += "Do you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
                buttons = [
                    {'text': 'Yes', 'value': 'yes', 'type': 'confirmation_choice'},
                    {'text': 'No', 'value': 'no', 'type': 'confirmation_choice'}
                ]
                return self._create_response_with_choice_buttons(response_text, thread_id, buttons)
        else:
            # Holistic fallback: accept a single date and treat it as a same-day range
            single = self.timeoff_service.parse_date_input(message)
            if single:
                self.session_manager.update_session(thread_id, {'start_date': single, 'end_date': single})
                self.session_manager.advance_session_step(thread_id)

                session_data = session.get('data', {})
                selected_type = session_data.get('selected_leave_type') or session.get('selected_leave_type', {})
                resolved_employee = employee_data or session_data.get('employee_data') or session.get('employee_data') or {}

                def dd_mm_yyyy(d: str) -> str:
                    try:
                        return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
                    except Exception:
                        return d
                if is_halfday_flow:
                    hour_text = (
                        "Great, got your date. Please choose your hours (from/to)."
                    )
                    return self._create_response_with_hour_picker(hour_text, thread_id)
                else:
                    response_text = f"Great, noted your date. Here's your time-off request summary:\n\n"
                    response_text += f"📋 **Leave Type:** {selected_type.get('name', 'Unknown') if selected_type else 'Unknown'}\n"
                    response_text += f"📅 **Start Date:** {dd_mm_yyyy(single)}\n"
                    response_text += f"📅 **End Date:** {dd_mm_yyyy(single)}\n"
                    response_text += f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}\n\n"
                    response_text += "Do you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
                    buttons = [
                        {'text': 'Yes', 'value': 'yes', 'type': 'confirmation_choice'},
                        {'text': 'No', 'value': 'no', 'type': 'confirmation_choice'}
                    ]
                    return self._create_response_with_choice_buttons(response_text, thread_id, buttons)

            response_text = (
                "I couldn't parse the date range. Please send both dates in one message. Examples:\n"
                "- 23/9 to 24/9\n"
                "- 23/09/2025 to 24/09/2025\n"
                "- 23-9-2025 till 24-9-2025\n"
                "- 23rd of September till the 24th\n"
                "- next Monday to Wednesday\n\n"
                "Defaults: I assume DD/MM, current month and year unless you specify otherwise."
            )
            return self._create_response(response_text, thread_id)
    
    def _handle_confirmation(self, message: str, thread_id: str, session: dict, employee_data: dict) -> dict:
        """Handle final confirmation step"""
        message_lower = message.lower().strip()
        debug_log(f"Handling confirmation - message: '{message_lower}'", "bot_logic")

        if message_lower in ['yes', 'y', 'confirm', 'submit', 'ok', 'sure']:
            # Submit the request
            debug_log(f"User confirmed submission, calling _submit_timeoff_request", "bot_logic")
            return self._submit_timeoff_request(thread_id, session, employee_data)
        elif message_lower in ['no', 'n', 'cancel', 'abort', 'stop', 'exit', 'quit', 'nevermind', 'end', 'undo']:
            # Cancel the request with standardized message
            debug_log(f"User cancelled submission", "bot_logic")
            try:
                self.session_manager.cancel_session(thread_id, 'User cancelled at confirmation')
            finally:
                self.session_manager.clear_session(thread_id)
            response_text = 'request cancelled, can i help you with anything else'
            return self._create_response(response_text, thread_id)
        else:
            # If we receive hour_from/hour_to structured message OR a natural hour range, capture and show summary
            try:
                parsed_from, parsed_to = self._parse_hour_range_text(message)
                has_structured = ('hour_from=' in message and 'hour_to=' in message)
                if has_structured or (parsed_from and parsed_to):
                    def _get_param(k: str, s: str) -> str:
                        try:
                            parts = {p.split('=')[0]: p.split('=')[1] for p in s.split('&') if '=' in p}
                            return parts.get(k, '')
                        except Exception:
                            return ''
                    raw_from = _get_param('hour_from', message) if has_structured else parsed_from
                    raw_to = _get_param('hour_to', message) if has_structured else parsed_to

                    session_data = session.get('data', {})
                    self.session_manager.update_session(thread_id, {'hour_from': raw_from, 'hour_to': raw_to})

                    selected_type = session_data.get('selected_leave_type') or session.get('selected_leave_type', {})
                    start_date = session_data.get('start_date') or session.get('start_date')
                    end_date = session_data.get('end_date') or session.get('end_date')
                    resolved_employee = employee_data or session_data.get('employee_data') or session.get('employee_data') or {}

                    def dd_mm_yyyy(d: str) -> str:
                        try:
                            return datetime.strptime(d, '%Y-%m-%d').strftime('%d/%m/%Y')
                        except Exception:
                            return d

                    # Half Day policy: max 4 hours
                    try:
                        hf = self._hour_key_to_float(raw_from)
                        ht = self._hour_key_to_float(raw_to)
                        if hf == hf and ht == ht and (ht - hf) > 4.0 - 1e-9:
                            policy_text = (
                                "Per Prezlab policy, the maximum Half Day leave duration is 4 hours. "
                                "Please enter hours less than or equal to 4."
                            )
                            return self._create_response_with_hour_picker(policy_text, thread_id)
                    except Exception:
                        pass

                    response_text = f"Great, noted your hours. Here's your time-off request summary:\n\n"
                    response_text += f"📋 **Leave Type:** {selected_type.get('name', 'Custom Hours')}\n"
                    response_text += f"📅 **Date:** {dd_mm_yyyy(start_date)}\n"
                    response_text += f"⏰ **Hours:** from {self._format_hour_label(raw_from)} to {self._format_hour_label(raw_to)}\n"
                    response_text += f"👤 **Employee:** {resolved_employee.get('name', 'Unknown')}\n\n"
                    response_text += "Do you want to submit this request? reply or click 'yes' to confirm or 'no' to cancel"
                    buttons = [
                        {'text': 'Yes', 'value': 'yes', 'type': 'confirmation_choice'},
                        {'text': 'No', 'value': 'no', 'type': 'confirmation_choice'}
                    ]
                    return self._create_response_with_choice_buttons(response_text, thread_id, buttons)
            except Exception:
                pass

            # Half Day UX: if invalid hour range (e.g., same start/end) or unparsed text, re-open the hour picker
            try:
                if self._is_halfday_flow(session):
                    hour_text = (
                        "Please choose a valid hours range (end must be after start)."
                    )
                    return self._create_response_with_hour_picker(hour_text, thread_id)
            except Exception:
                pass

            debug_log(f"Invalid confirmation response: '{message_lower}'", "bot_logic")
            response_text = "Please reply with 'yes' to submit the request or 'no' to cancel."
            return self._create_response(response_text, thread_id)
    
    def _submit_timeoff_request(self, thread_id: str, session: dict, employee_data: dict) -> dict:
        """Submit the time-off request to Odoo"""
        try:
            debug_log(f"Starting time-off submission for thread: {thread_id}", "bot_logic")
            debug_log(f"Full session structure: {session}", "bot_logic")

            session_data = session.get('data', {})
            # Check both session root level and nested data level
            selected_type = session_data.get('selected_leave_type') or session.get('selected_leave_type', {})
            start_date = session_data.get('start_date') or session.get('start_date')
            end_date = session_data.get('end_date') or session.get('end_date')
            # Resolve employee data from argument or session to ensure correctness
            resolved_employee = employee_data or session_data.get('employee_data') or session.get('employee_data') or {}

            debug_log(f"Session data - selected_type: {selected_type}, start_date: {start_date}, end_date: {end_date}", "bot_logic")

            employee_id = resolved_employee.get('id') if resolved_employee else None
            leave_type_id = selected_type.get('id') if selected_type else None

            debug_log(f"Extracted data - employee_id: {employee_id}, leave_type_id: {leave_type_id}", "bot_logic")

            if not all([employee_id, leave_type_id, start_date, end_date]):
                debug_log(f"Missing required data - employee_id: {employee_id}, leave_type_id: {leave_type_id}, start_date: {start_date}, end_date: {end_date}", "general")
                raise ValueError("Missing required data for submission")
            
            # Support Half Day custom hours via modular service and hour range fields
            extra_fields = {}
            if self.halfday_service and isinstance(selected_type, dict):
                try:
                    mapped_leave_type_id, hd_extra = self.halfday_service.build_submission(selected_type)
                    if mapped_leave_type_id:
                        leave_type_id = mapped_leave_type_id
                    if hd_extra:
                        extra_fields.update(hd_extra)
                except Exception as hd_map_e:
                    debug_log(f"HalfDay build_submission error: {hd_map_e}", "general")
                # Add hour range if present in session
                sd = session.get('data', {})
                hour_from = sd.get('hour_from') or session.get('hour_from')
                hour_to = sd.get('hour_to') or session.get('hour_to')

                # Enforce 4-hour maximum
                try:
                    hf = float(hour_from) if hour_from is not None else float('nan')
                    ht = float(hour_to) if hour_to is not None else float('nan')
                    if hf == hf and ht == ht and (ht - hf) > 4.0 - 1e-9:
                        return self._create_response_with_hour_picker(
                            "Per Prezlab policy, the maximum Half Day leave duration is 4 hours. Please choose up to 4 hours.",
                            thread_id
                        )
                except Exception:
                    pass

                def _to_selection_key(val):
                    # Odoo selection keys are strings like '16' or '16.5'
                    try:
                        f = float(val) if isinstance(val, (str, int, float)) else None
                        if f is None:
                            return None
                        return str(int(f)) if abs(f - int(f)) < 1e-9 else str(f)
                    except Exception:
                        return None

                hf = _to_selection_key(hour_from)
                ht = _to_selection_key(hour_to)
                if hf is not None:
                    extra_fields['request_hour_from'] = hf
                if ht is not None:
                    extra_fields['request_hour_to'] = ht

            success, result = self.timeoff_service.submit_leave_request(
                employee_id=employee_id,
                leave_type_id=leave_type_id,
                start_date=start_date,
                end_date=end_date,
                description=f"Time off request submitted via Nasma chatbot",
                extra_fields=extra_fields or None
            )
            
            if success:
                self.session_manager.complete_session(thread_id, {'submitted': True, 'result': result})
                # Immediately clear session so a new flow can start right away
                try:
                    self.session_manager.clear_session(thread_id)
                except Exception:
                    pass
                response_text = f"✅ **Success!** {result.get('message', 'Your time-off request has been submitted.')}\n\n"
                response_text += "Your request is now pending approval from your manager. You should receive a notification once it's reviewed."
            else:
                self.session_manager.complete_session(thread_id, {'submitted': False, 'error': result})
                try:
                    self.session_manager.clear_session(thread_id)
                except Exception:
                    pass
                response_text = f"❌ **Submission Failed:** {result}\n\n"
                response_text += "Please try again later or contact your HR department for assistance."
            
            resp = self._create_response(response_text, thread_id)
            try:
                # Mark as handled to avoid any lingering prompts
                resp['session_handled'] = True
            except Exception:
                pass
            return resp
            
        except Exception as e:
            debug_log(f"Error submitting time-off request: {e}", "general")
            self.session_manager.complete_session(thread_id, {'submitted': False, 'error': str(e)})
            response_text = "❌ Sorry, there was an error submitting your request. Please try again later or contact HR for assistance."
            return self._create_response(response_text, thread_id)
    
    def _create_response(self, message_text: str, thread_id: str) -> dict:
        """Create a standard response object"""
        debug_log(f"Creating response with message length: {len(message_text) if message_text else 0}", "bot_logic")
        debug_log(f"Response thread_id: {thread_id}", "bot_logic")

        # Ensure we always have a thread_id for session continuity
        if not thread_id:
            import time
            thread_id = f"timeoff_{int(time.time())}"
            debug_log(f"Generated fallback thread_id in response: {thread_id}", "bot_logic")

        result = {
            'message': message_text,
            'thread_id': thread_id,
            'source': self.model,
            'confidence_score': 1.0,
            'model_used': self.model,
            'session_handled': True
        }

        debug_log(f"Response object created successfully: {result is not None}", "bot_logic")
        return result

    def _create_response_with_buttons(self, message_text: str, thread_id: str, leave_types: list) -> dict:
        """Create a response object with interactive buttons"""
        debug_log(f"Creating response with buttons - message length: {len(message_text) if message_text else 0}", "bot_logic")
        debug_log(f"Response thread_id: {thread_id}", "bot_logic")

        # Ensure we always have a thread_id for session continuity
        if not thread_id:
            import time
            thread_id = f"timeoff_{int(time.time())}"
            debug_log(f"Generated fallback thread_id in response with buttons: {thread_id}", "bot_logic")

        # Create button data
        buttons = []
        for leave_type in leave_types:
            buttons.append({
                'text': leave_type.get('name', 'Unknown'),
                'value': leave_type.get('name', 'Unknown'),
                'type': 'leave_type_selection'
            })

        result = {
            'message': message_text,
            'thread_id': thread_id,
            'source': self.model,
            'confidence_score': 1.0,
            'model_used': self.model,
            'session_handled': True,
            'buttons': buttons  # Add buttons to response
        }

        debug_log(f"Response object with {len(buttons)} buttons created successfully", "bot_logic")
        return result

    def _create_response_with_datepicker(self, message_text: str, thread_id: str) -> dict:
        """Create a response object that instructs the UI to show a date range picker widget"""
        debug_log(f"Creating response with datepicker - message length: {len(message_text) if message_text else 0}", "bot_logic")
        if not thread_id:
            import time
            thread_id = f"timeoff_{int(time.time())}"
            debug_log(f"Generated fallback thread_id in response with datepicker: {thread_id}", "bot_logic")

        return {
            'message': message_text,
            'thread_id': thread_id,
            'source': self.model,
            'confidence_score': 1.0,
            'model_used': self.model,
            'session_handled': True,
            'widgets': {
                'date_range_picker': True
            }
        }

    def _create_response_with_datepicker_single(self, message_text: str, thread_id: str) -> dict:
        """Create a response object to show a single-date picker widget (for Half Days)."""
        debug_log(f"Creating response with single-date picker - message length: {len(message_text) if message_text else 0}", "bot_logic")
        if not thread_id:
            import time
            thread_id = f"timeoff_{int(time.time())}"
            debug_log(f"Generated fallback thread_id in response with single datepicker: {thread_id}", "bot_logic")

        return {
            'message': message_text,
            'thread_id': thread_id,
            'source': self.model,
            'confidence_score': 1.0,
            'model_used': self.model,
            'session_handled': True,
            'widgets': {
                'single_date_picker': True
            }
        }

    def _create_response_with_hour_picker(self, message_text: str, thread_id: str) -> dict:
        """Create a response object that instructs the UI to show an hour-from/to picker."""
        debug_log(f"Creating response with hour picker - message length: {len(message_text) if message_text else 0}", "bot_logic")
        if not thread_id:
            import time
            thread_id = f"timeoff_{int(time.time())}"
            debug_log(f"Generated fallback thread_id in response with hour picker: {thread_id}", "bot_logic")

        # Build hour options (0..23.5 step 0.5)
        def _fmt_label(val: float) -> str:
            h = int(val)
            m = 30 if abs(val - h - 0.5) < 1e-6 else 0
            ampm_h = h % 12
            if ampm_h == 0:
                ampm_h = 12
            ampm = 'AM' if h < 12 else 'PM'
            return f"{ampm_h}:{m:02d} {ampm}" if m else f"{ampm_h}:00 {ampm}"

        options = []
        def _push(val: float):
            key = str(int(val)) if abs(val - int(val)) < 1e-9 else str(val)
            options.append({'value': key, 'label': _fmt_label(val)})

        # From 9:00 to 23:30 (same day)
        v = 9.0
        while v <= 23.5 + 1e-9:
            _push(v)
            v += 0.5
        # Then wrap 0:00 to 1:00 (next day)
        v = 0.0
        while v <= 1.0 + 1e-9:
            _push(v)
            v += 0.5

        return {
            'message': message_text,
            'thread_id': thread_id,
            'source': self.model,
            'confidence_score': 1.0,
            'model_used': self.model,
            'session_handled': True,
            'widgets': {
                'hour_range_picker': True,
                'hour_options': options
            }
        }

    def _is_halfday_flow(self, session: dict) -> bool:
        session_data = session.get('data', {}) if isinstance(session, dict) else {}
        selected_type = session_data.get('selected_leave_type') or session.get('selected_leave_type', {})
        try:
            if self.halfday_service and isinstance(selected_type, dict):
                return self.halfday_service.is_halfday(selected_type) or (
                    selected_type.get('name') == getattr(self.halfday_service, 'HALF_DAY_NAME', 'Custom Hours')
                )
        except Exception:
            return False
        return False

    def _format_hour_label(self, value: str) -> str:
        """Format a float-like string hour value to 12-hour label (e.g., '8.5' -> '8:30 AM')."""
        try:
            f = float(value)
            h = int(f)
            minutes = 30 if abs(f - h - 0.5) < 1e-6 else 0
            ampm_h = h % 12
            if ampm_h == 0:
                ampm_h = 12
            ampm = 'AM' if h < 12 else 'PM'
            return f"{ampm_h}:{minutes:02d} {ampm}" if minutes else f"{ampm_h}:00 {ampm}"
        except Exception:
            return value

    def _create_response_with_choice_buttons(self, message_text: str, thread_id: str,
                                             buttons: list) -> dict:
        """Create a response with explicit buttons, each a dict with text/value/type."""
        debug_log(f"Creating response with choice buttons - count: {len(buttons) if buttons else 0}", "bot_logic")
        if not thread_id:
            import time
            thread_id = f"timeoff_{int(time.time())}"
        return {
            'message': message_text,
            'thread_id': thread_id,
            'source': self.model,
            'confidence_score': 1.0,
            'model_used': self.model,
            'session_handled': True,
            'buttons': buttons
        }

    def _parse_hour_value(self, token: str) -> float:
        """Parse a time token like '4pm', '16', '4:30 pm', '16:30' to 24h float (e.g., 16.5)."""
        try:
            import re
            s = token.strip().lower()
            s = s.replace('.', ':')
            m = re.match(r"^(\d{1,2})(?::(\d{2}))?\s*(am|pm)?$", s)
            if not m:
                return float('nan')
            h = int(m.group(1))
            mins = int(m.group(2)) if m.group(2) else 0
            ap = m.group(3)
            if ap == 'am':
                if h == 12:
                    h = 0
            elif ap == 'pm':
                if h != 12:
                    h += 12
            # clamp to 0..23:59
            h = max(0, min(23, h))
            mins = 0 if mins < 15 else (30 if mins < 45 else 0)  # snap to nearest half hour (0 or 30)
            return h + (0.5 if mins == 30 else 0.0)
        except Exception:
            return float('nan')

    def _parse_hour_range_text(self, text: str):
        """Parse expressions like '4pm to 5pm', '16:00 - 17:30'. Return (from_str, to_str) or (None, None)."""
        try:
            import re
            if not text or not isinstance(text, str):
                return None, None
            s = text.strip().lower()
            # Normalize connectors
            s = re.sub(r"\s*(?:–|—)\s*", '-', s)
            parts = re.split(r"\s*(?:to|till|until|-)\s*", s)
            if len(parts) != 2:
                return None, None
            v1 = self._parse_hour_value(parts[0])
            v2 = self._parse_hour_value(parts[1])
            if v1 != v1 or v2 != v2:  # NaN check
                return None, None
            # Ensure within 0..24 and logical order
            if v2 <= v1:
                return None, None
            return str(v1), str(v2)
        except Exception:
            return None, None

    def _hour_key_to_float(self, val: str) -> float:
        try:
            return float(val)
        except Exception:
            return float('nan')